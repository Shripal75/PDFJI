import os
from pypdf import PdfReader, PdfWriter

import fitz  # PyMuPDF
from PIL import Image
import io

def merge_pdfs(pdf_paths, output_path):
    writer = PdfWriter()
    for pdf in pdf_paths:
        reader = PdfReader(pdf)
        for page in reader.pages:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

import fitz  # PyMuPDF
from PIL import Image
import io

def compress_pdf(input_path, output_path, quality=50, target_size=None):
    """
    Compress PDF by re-compressing all images at the given JPEG quality.
    quality: int (0-100), where lower = smaller file, higher = better image quality.
    If target_size is provided (in bytes), it will try to hit that exact size.
    """
    
    def run_compression(q, out_path, scale=1.0):
        """Compress PDF images. q = JPEG quality (1-95), scale = image dimension scale (0.1-1.0)."""
        try:
            doc = fitz.open(input_path)
            doc.scrub(
                clean_pages=True,
                embedded_files=True,
                hidden_text=False,
                javascript=True,
                metadata=True,
                thumbnails=True,
                xml_metadata=True
            )
            
            for page in doc:
                image_list = page.get_images(full=True)
                for img_info in image_list:
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        if not base_image: continue
                        
                        img = Image.open(io.BytesIO(base_image["image"]))
                        
                        if q < 15:
                            img = img.convert("L")
                        elif img.mode in ("RGBA", "P", "LA"):
                            img = img.convert("RGB")
                        
                        # Only resize if scale < 1.0 (second phase)
                        if scale < 1.0:
                            new_size = (max(1, int(img.width * scale)), max(1, int(img.height * scale)))
                            img = img.resize(new_size, Image.Resampling.LANCZOS)
                        
                        buf = io.BytesIO()
                        img.save(buf, format="JPEG", quality=max(1, q), optimize=True)
                        compressed_bytes = buf.getvalue()
                        
                        if len(compressed_bytes) < len(base_image["image"]):
                            doc.update_stream(xref, compressed_bytes)
                            doc.xref_set_key(xref, "Filter", "/DCTDecode")
                    except: continue
            
            try:
                doc.subset_fonts()
            except:
                pass
                
            doc.save(out_path, garbage=4, deflate=True, clean=True)
            doc.close()
            return os.path.getsize(out_path)
        except Exception as e:
            print(f"Pass failed: {e}")
            return os.path.getsize(input_path)

    def run_rasterization(q, out_path, dpi=150):
        """Convert every page to an image and rebuild the PDF (Nuclear Option)."""
        try:
            doc = fitz.open(input_path)
            images = []
            matrix = fitz.Matrix(dpi / 72, dpi / 72)
            
            for page in doc:
                pix = page.get_pixmap(matrix=matrix, colorspace=fitz.csRGB if q >= 15 else fitz.csGRAY)
                img = Image.open(io.BytesIO(pix.tobytes("jpeg")))
                images.append(img)
            
            if images:
                images[0].save(
                    out_path, "PDF", resolution=float(dpi), save_all=True, append_images=images[1:],
                    quality=max(1, q), optimize=True
                )
            
            doc.close()
            return os.path.getsize(out_path)
        except Exception as e:
            print(f"Rasterization failed: {e}")
            return os.path.getsize(input_path)

    if target_size:
        # Phase 1: Binary search on JPEG quality only (no resizing)
        low = 1
        high = 95
        best_q = low
        best_size = None
        
        for _ in range(10):
            mid = (low + high) // 2
            current_size = run_compression(mid, output_path)
            
            if current_size <= target_size:
                best_q = mid
                best_size = current_size
                low = mid + 1  # Try higher quality (closer to target)
            else:
                high = mid - 1
            
            if low > high: break
        
        current_size = run_compression(best_q, output_path)
        
        # Phase 2: If quality-only compression still over target, add scaling
        if current_size > target_size:
            scale_low = 0.1
            scale_high = 1.0
            best_scale = scale_low
            
            for _ in range(8):
                scale_mid = (scale_low + scale_high) / 2
                current_size = run_compression(1, output_path, scale=scale_mid)
                
                if current_size <= target_size:
                    best_scale = scale_mid
                    scale_low = scale_mid  # Try larger scale (closer to target)
                else:
                    scale_high = scale_mid
                
                if scale_high - scale_low < 0.02: break
            
            current_size = run_compression(1, output_path, scale=best_scale)
        
        # Phase 3: Nuclear — rasterize if still over target
        if current_size > target_size:
            dpi_low = 36
            dpi_high = 150
            best_dpi = dpi_low
            
            for _ in range(6):
                dpi_mid = (dpi_low + dpi_high) // 2
                current_size = run_rasterization(10, output_path, dpi=dpi_mid)
                
                if current_size <= target_size:
                    best_dpi = dpi_mid
                    dpi_low = dpi_mid + 1
                else:
                    dpi_high = dpi_mid - 1
                
                if dpi_low > dpi_high: break
            
            run_rasterization(10, output_path, dpi=best_dpi)
    else:
        run_compression(quality, output_path)

    return output_path

def convert_pdf_to_word(input_path, output_path, pages=None):
    from pdf2docx import Converter
    cv = Converter(input_path)
    if pages:
        cv.convert(output_path, pages=pages)
    else:
        cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path

def images_to_pdf(image_paths, output_path):
    """
    Convert a list of images to a single PDF.
    image_paths: List of file paths to images.
    """
    images = []
    try:
        for path in image_paths:
            img = Image.open(path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            images.append(img)
        
        if images:
            images[0].save(
                output_path, "PDF", resolution=100.0, save_all=True, append_images=images[1:]
            )
    except Exception as e:
        print(f"Image to PDF conversion failed: {e}")
        raise e
    return output_path

def pdf_to_images(pdf_path, output_dir):
    """
    Convert PDF pages to images for preview.
    Returns list of image filenames.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    doc = fitz.open(pdf_path)
    image_paths = []
    
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5)) # Low res for preview
        image_name = f"page_{i}.png"
        image_path = os.path.join(output_dir, image_name)
        pix.save(image_path)
        image_paths.append(image_name) # Return relative filename
        
    return image_paths

def pdf_to_high_res_images(pdf_path, output_dir, fmt="png", page_indices=None):
    """
    Convert PDF pages to high-resolution images for download.
    fmt: 'png' or 'jpg'
    page_indices: optional list of 0-based page indices. If None, all pages.
    Returns list of image file paths.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    doc = fitz.open(pdf_path)
    image_paths = []
    
    indices = page_indices if page_indices is not None else list(range(len(doc)))
    
    for i in indices:
        if 0 <= i < len(doc):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            ext = "jpg" if fmt == "jpg" else "png"
            image_name = f"page_{i + 1}.{ext}"
            image_path = os.path.join(output_dir, image_name)
            if fmt == "jpg":
                pix.save(image_path, jpg_quality=95)
            else:
                pix.save(image_path)
            image_paths.append(image_path)
    
    doc.close()
    return image_paths

def reorder_pdf(pdf_path, output_path, page_order):
    """
    Create a new PDF with pages in the specified order.
    page_order: List of integer indices (0-based).
    """
    doc = fitz.open(pdf_path)
    
    # page_order can contain duplicates if user duplicating pages? 
    # Or just subset.
    valid_indices = [i for i in page_order if 0 <= i < len(doc)]
    
    doc.select(valid_indices)
    doc.save(output_path)
    doc.close()
    
    return output_path

def parse_page_range(range_str, max_pages):
    """
    Parse range string like "1-3, 5" into 0-based index list.
    """
    pages = set()
    parts = range_str.split(',')
    for part in parts:
        part = part.strip()
        if '-' in part:
            start, end = part.split('-')
            try:
                start = int(start)
                end = int(end)
                for i in range(start, end + 1):
                    if 1 <= i <= max_pages:
                        pages.add(i - 1)
            except ValueError:
                continue
        else:
            try:
                page = int(part)
                if 1 <= page <= max_pages:
                    pages.add(page - 1)
            except ValueError:
                continue
    return sorted(list(pages))

def split_pdf(input_path, output_path, page_range):
    """
    Create a new PDF with selected pages.
    """
    doc = fitz.open(input_path)
    pages = parse_page_range(page_range, len(doc))
    
    if not pages:
        raise ValueError("No valid pages selected")
        
    doc.select(pages)
    doc.save(output_path)
    doc.close()
    return output_path

def rotate_pdf(input_path, output_path, rotations):
    """
    Rotate individual pages.
    rotations: dict mapping page_index (0-based) to angle (90, 180, 270).
    Example: {0: 90, 2: 180} rotates page 1 by 90° and page 3 by 180°.
    """
    doc = fitz.open(input_path)
    for page_idx_str, angle in rotations.items():
        page_idx = int(page_idx_str)
        if 0 <= page_idx < len(doc):
            page = doc[page_idx]
            page.set_rotation(page.rotation + angle)
        
    doc.save(output_path)
    doc.close()
    return output_path

def extract_text(input_path, output_path, page_indices=None):
    """
    Extract text content to a .txt file.
    page_indices: optional list of 0-based page indices.
    """
    doc = fitz.open(input_path)
    full_text = ""
    
    indices = page_indices if page_indices is not None else range(len(doc))
    
    for i in indices:
        if 0 <= i < len(doc):
            page = doc.load_page(i)
            text = page.get_text()
            if text:
               full_text += f"--- Page {i+1} ---\n{text}\n\n"
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(full_text)
    
    doc.close()
    return output_path

def word_to_pdf(input_path, output_path):
    """
    Convert Word document (.docx) to PDF.
    Tries docx2pdf (needs MS Word) first, falls back to pure Python.
    """
    # Method 1: Try docx2pdf (requires MS Word installed)
    try:
        from docx2pdf import convert as docx2pdf_convert
        docx2pdf_convert(input_path, output_path)
        return output_path
    except Exception:
        print("docx2pdf failed (MS Word not available), using pure Python fallback...")

    # Method 2: Pure Python using python-docx + PyMuPDF
    from docx import Document
    from docx.shared import Pt
    
    doc = Document(input_path)
    pdf = fitz.open()
    
    # A4 dimensions in points
    page_w, page_h = 595.28, 841.89
    margin = 72  # 1 inch
    usable_w = page_w - 2 * margin
    
    page = pdf.new_page(width=page_w, height=page_h)
    y = margin
    
    def new_page():
        nonlocal page, y
        page = pdf.new_page(width=page_w, height=page_h)
        y = margin
    
    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            y += 8
            if y > page_h - margin:
                new_page()
            continue
        
        # Determine font size from style
        font_size = 11
        fontname = "helv"
        style_name = para.style.name.lower() if para.style else ""
        
        if "heading 1" in style_name:
            font_size = 22
            fontname = "hebo"
        elif "heading 2" in style_name:
            font_size = 17
            fontname = "hebo"
        elif "heading 3" in style_name:
            font_size = 14
            fontname = "hebo"
        elif "title" in style_name:
            font_size = 26
            fontname = "hebo"
        
        # Check if any run is bold
        if any(run.bold for run in para.runs if run.bold is not None):
            if fontname == "helv":
                fontname = "hebo"
        
        line_height = font_size + 4
        
        # Estimate how many lines this text will need
        chars_per_line = max(int(usable_w / (font_size * 0.5)), 20)
        estimated_lines = max(1, len(text) // chars_per_line + 1)
        text_block_height = estimated_lines * line_height + 4
        
        # Need new page?
        if y + text_block_height > page_h - margin:
            new_page()
        
        # Insert text with wrapping
        rect = fitz.Rect(margin, y, page_w - margin, page_h - margin)
        rc = page.insert_textbox(
            rect, text,
            fontsize=font_size,
            fontname=fontname,
            align=0  # left align
        )
        
        # rc is the unused height (negative means overflow)
        if rc < 0:
            # Text overflowed, put remainder on new page
            new_page()
            rect = fitz.Rect(margin, y, page_w - margin, page_h - margin)
            page.insert_textbox(rect, text, fontsize=font_size, fontname=fontname)
        
        y += text_block_height
    
    pdf.save(output_path)
    pdf.close()
    return output_path

def compress_image(input_path, output_path, quality=50, target_size=None):
    """
    Compress image files (JPG, PNG, WEBP).
    If target_size is provided (in bytes), it tries to hit that exact size.
    """
    img = Image.open(input_path)
    original_format = img.format
    
    # Handle transparency for formats that don't support it
    if img.mode in ("RGBA", "P", "LA") and original_format == "JPEG":
        img = img.convert("RGB")

    def run_compression(q, out_path, scale=1.0):
        """Compress image. q = quality (1-95), scale = dimension scale (0.1-1.0)."""
        try:
            comp_img = Image.open(input_path)
            
            if comp_img.mode in ("RGBA", "P", "LA") and (original_format == "JPEG" or q < 10):
                comp_img = comp_img.convert("RGB")
            
            if q < 10:
                comp_img = comp_img.convert("L")

            # Only resize if scale < 1.0
            if scale < 1.0:
                new_size = (max(1, int(comp_img.width * scale)), max(1, int(comp_img.height * scale)))
                comp_img = comp_img.resize(new_size, Image.Resampling.LANCZOS)

            save_args = {"optimize": True}
            
            if original_format == "PNG":
                save_args["format"] = "PNG"
            elif original_format == "WEBP":
                save_args["format"] = "WEBP"
                save_args["quality"] = max(1, q)
            else:
                save_args["format"] = "JPEG"
                save_args["quality"] = max(1, q)

            buf = io.BytesIO()
            comp_img.save(buf, **save_args)
            compressed_bytes = buf.getvalue()
            
            with open(out_path, "wb") as f:
                f.write(compressed_bytes)
                
            return len(compressed_bytes)
        except Exception as e:
            print(f"Image compression pass failed: {e}")
            return os.path.getsize(input_path)

    if target_size:
        # Phase 1: Binary search on quality only (no resizing)
        low = 1
        high = 95
        best_q = low
        
        for _ in range(10):
            mid = (low + high) // 2
            current_size = run_compression(mid, output_path)
            
            if current_size <= target_size:
                best_q = mid
                low = mid + 1  # Try higher quality (closer to target)
            else:
                high = mid - 1
            
            if low > high: break
        
        current_size = run_compression(best_q, output_path)
        
        # Phase 2: If quality alone can't reach target, add scaling
        if current_size > target_size:
            scale_low = 0.1
            scale_high = 1.0
            best_scale = scale_low
            
            for _ in range(10):
                scale_mid = (scale_low + scale_high) / 2
                current_size = run_compression(1, output_path, scale=scale_mid)
                
                if current_size <= target_size:
                    best_scale = scale_mid
                    scale_low = scale_mid  # Try larger scale (closer to target)
                else:
                    scale_high = scale_mid
                
                if scale_high - scale_low < 0.01: break
            
            run_compression(1, output_path, scale=best_scale)
    else:
        # Simple quality-only compression if no target size
        img = Image.open(input_path)
        if img.mode in ("RGBA", "P", "LA") and original_format == "JPEG":
            img = img.convert("RGB")
        
        save_args = {"optimize": True}
        if original_format == "PNG":
            save_args["format"] = "PNG"
        elif original_format == "WEBP":
             save_args["format"] = "WEBP"
             save_args["quality"] = quality
        else:
             save_args["format"] = "JPEG"
             save_args["quality"] = quality
             
        img.save(output_path, **save_args)
        
    return output_path

def pdf_to_ppt(input_path, output_path, page_indices=None):
    """
    Convert PDF to PPTX by rendering each page as an image and adding to a slide.
    page_indices: optional list of 0-based page indices to convert.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        import io 
    except ImportError:
        raise ImportError("python-pptx is required. pip install python-pptx")

    doc = fitz.open(input_path)
    prs = Presentation()
    
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6] 
    
    # Determine pages to process
    indices = page_indices if page_indices is not None else range(len(doc))
    
    for i in indices:
        if 0 <= i < len(doc):
            page = doc.load_page(i)
            
            # Render high-res image
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            img_data = pix.tobytes("png")
            image_stream = io.BytesIO(img_data)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide, filling it
            slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save(output_path)
    doc.close()
    return output_path

def ppt_to_pdf(input_path, output_path, page_indices=None):
    """
    Convert PPTX to PDF using Microsoft PowerPoint via COM (Windows only).
    page_indices: optional list of 0-based page indices to convert.
    """
    try:
        import comtypes.client
    except ImportError:
        raise ImportError("comtypes is required for PPT to PDF (Windows only). pip install comtypes")
    
    # COM requires absolute paths
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
        
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    
    try:
        # Open presentation
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        
        # Filter pages if requested
        if page_indices is not None:
            # Create set for O(1) lookup
            keep_indices = set(page_indices)
            count = presentation.Slides.Count
            
            # Delete slides backwards to avoid index shifting
            for i in range(count, 0, -1):
                # i is 1-based, so i-1 is 0-based index
                if (i - 1) not in keep_indices:
                    presentation.Slides(i).Delete()
        
        # Save as PDF (format 32)
        presentation.SaveAs(output_path, 32)
        
        presentation.Close()
    except Exception as e:
        raise e
    finally:
        try:
            powerpoint.Quit()
        except:
            pass
    
    return output_path

def extract_pptx_slides(input_path, output_dir):
    """
    Extract all slides from PPTX as images using COM.
    Returns list of image filenames.
    """
    try:
        import comtypes.client
    except ImportError:
        return []
    
    # COM requires absolute paths
    input_path = os.path.abspath(input_path)
    output_dir = os.path.abspath(output_dir)
        
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    
    image_paths = []
    
    try:
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        
        for i, slide in enumerate(presentation.Slides):
            # 1-based index for naming
            image_name = f"slide_{i+1}.jpg"
            image_path = os.path.join(output_dir, image_name)
            
            # Export slide
            slide.Export(image_path, "JPG")
            image_paths.append(image_name)
            
        presentation.Close()
    except Exception as e:
        print(f"Error extracting PPTX slides: {e}")
    finally:
        try:
            powerpoint.Quit()
        except:
            pass
            
    return image_paths

def assemble_pdf(pages_spec, output_path):
    """
    Assemble a PDF from multiple sources with rotation.
    pages_spec: list of dicts {'file_path': str, 'page_index': int, 'rotate': int}
    """
    out_doc = fitz.open()
    
    # Cache open documents to avoid re-opening for every page
    docs_cache = {}
    
    try:
        for spec in pages_spec:
            fp = spec['file_path']
            if fp not in docs_cache:
                docs_cache[fp] = fitz.open(fp)
            
            src_doc = docs_cache[fp]
            p_index = spec['page_index']
            rotation = spec.get('rotate', 0)
            
            if 0 <= p_index < len(src_doc):
                out_doc.insert_pdf(src_doc, from_page=p_index, to_page=p_index, rotate=rotation)
                
        out_doc.save(output_path)
    finally:
        for doc in docs_cache.values():
            doc.close()
        out_doc.close()
    
    return output_path


def convert_image(input_path, output_path, target_format):
    """
    Convert an image to a different format using Pillow.
    target_format: 'png', 'jpg', or 'webp'
    """
    img = Image.open(input_path)

    # Convert RGBA to RGB for formats that don't support alpha
    if target_format.lower() in ('jpg', 'jpeg') and img.mode in ('RGBA', 'P', 'LA'):
        bg = Image.new('RGB', img.size, (255, 255, 255))
        if img.mode == 'P':
            img = img.convert('RGBA')
        bg.paste(img, mask=img.split()[-1] if 'A' in img.mode else None)
        img = bg

    save_format = 'JPEG' if target_format.lower() in ('jpg', 'jpeg') else target_format.upper()
    save_kwargs = {'quality': 95} if save_format == 'JPEG' else {}

    img.save(output_path, format=save_format, **save_kwargs)
    return output_path
